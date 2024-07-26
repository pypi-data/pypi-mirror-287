# split_txts.py
import argparse
from typing import Optional, List

import fitz  # PyMuPDF
import os
import pandas as pd
import re
import subprocess
from docx import Document
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from botrun_ask_folder.models.rag_metadata import RagMetadata
from .drive_download_metadata import update_download_metadata
from .emoji_progress_bar import EmojiProgressBar
from .is_cached import is_cached
from .is_text_file import is_text_file
from .log_handler import LogHandler
from .safe_join import safe_join

log_handler = LogHandler("./logs", "split_txts.py")  # 創建日誌處理器實例


def get_page_number_prefix():
    return "page_"  # f"page_{page_num}.txt"


def get_file_name_with_page_number(file_name: str, page_num: int) -> str:
    return f"{file_name}.{get_page_number_prefix()}{page_num}.txt"


def re_compile_page_number_pattern():
    return re.compile(rf'\.{get_page_number_prefix()}\d+\.txt$')


def get_page_number(str_file_path):
    pattern = re.compile(rf'\.{get_page_number_prefix()}(\d+)\.txt$')
    match = pattern.search(str_file_path)
    if match:
        return int(match.group(1))
    else:
        return 0


def write_to_file(output_folder: str, filename: str, text: str, chars_per_page: int, add_page_numbers: bool = True):
    output_path = Path(output_folder)
    # pdf like files, 如果 add_page_numbers 為 False，則不加入頁碼，並且完整寫入檔案不要切割，此時 chars_per_page 無用
    if not add_page_numbers:
        output_filename = f"{filename}.txt"
        with open(output_path / output_filename, 'w', encoding='utf-8') as file:
            file.write(text)
        return
    # txt like files, 需要切割, 按照 chars_per_page 大小切割
    for page_num, i in enumerate(range(0, len(text), chars_per_page), start=1):
        page_filename = get_file_name_with_page_number(filename, page_num)
        with open(output_path / page_filename, 'w', encoding='utf-8') as file:
            file.write(text[i:i + chars_per_page])


def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)
    return ' '.join(text)


def extract_text_from_pptx(slide):
    text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text.append(shape.text)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        text.append(cell.text_frame.text)
    return ' '.join(text)


def process_powerpoint(original_file_path, file_type, output_folder, chars_per_page):
    file_name = Path(original_file_path).stem
    expected_pptx_path = safe_join(output_folder, f'{file_name}.pptx')
    if not os.path.exists(expected_pptx_path) or file_type == '.ppt':
        if file_type == '.ppt':
            converted_file_path = convert_office_file(original_file_path, file_type)
        else:
            converted_file_path = original_file_path
    else:
        converted_file_path = expected_pptx_path
    prs = Presentation(converted_file_path)
    for i, slide in enumerate(prs.slides):
        text = extract_text_from_pptx(slide)
        write_to_file(output_folder, f'{file_name}.{get_page_number_prefix()}{i + 1}', text, chars_per_page, False)


def process_excel(original_file_path: str, output_folder: str, file_name: str, file_type: str):
    if file_type == '.xlsx':
        xls = pd.ExcelFile(original_file_path, engine='openpyxl')
    elif file_type == '.xls':
        xls = pd.ExcelFile(original_file_path, engine='xlrd')
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name)
        # todo 2024-03-10 bowen, 每個 sheet 目前只有切一頁沒有分頁，有機率太大而無法放入大模型
        csv_output_path = f"{output_folder}/{file_name}_{sheet_name}.{get_page_number_prefix()}1.txt"
        df.to_csv(csv_output_path, index=False, encoding='utf-8-sig', quoting=0)


def convert_office_file(file_path: str, target_extension: str):
    converted_file_dir = os.path.dirname(file_path)
    converted_file_base = os.path.splitext(os.path.basename(file_path))[0]
    converted_file_name = f"{converted_file_base}.{target_extension.lstrip('.')}"
    converted_file_path = safe_join(converted_file_dir, converted_file_name)

    # 建構 LibreOffice 的命令列格式，用於檔案轉換
    libreoffice_command = os.getenv("LIBRE_OFFICE_COMMAND", "libreoffice")
    cmd = [libreoffice_command, '--headless', '--convert-to', target_extension.lstrip('.'), '--outdir',
           converted_file_dir,
           file_path]
    try:
        proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=True)
        stdout, stderr = proc.stdout, proc.stderr
        log_handler.write_log(f"Stdout: {stdout}")
        log_handler.write_log(f"Stderr: {stderr}")
    except subprocess.CalledProcessError as e:
        error_message = f"Error converting {file_path}. Return code: {e.returncode}\nStdout: {e.stdout}\nStderr: {e.stderr}"
        log_handler.write_log(error_message)
        raise Exception(error_message)
    except Exception as e:
        log_handler.write_log(f"Error converting {file_path}: {e}")
        raise Exception(f"Error converting {file_path}: {e}")
    if proc.returncode != 0:
        log_handler.write_log(f"Error converting {file_path}: {stderr}")
        raise Exception(f"Error converting {file_path}: {stderr}")
    return converted_file_path


def process_single_file(source_file: str,
                        output_folder: str,
                        chars_per_page: int,
                        force: bool,
                        stats: dict,
                        gen_page_imgs: bool = False, ) -> List[RagMetadata]:
    """
    @param metadata_dir: str, metadata 資料夾路徑，有傳才會更新 metadata
    """
    file_name, file_type = os.path.splitext(os.path.basename(source_file))
    if is_cached(output_folder, force):
        stats['cached'] += 1
        return []
    os.makedirs(output_folder, exist_ok=True)
    initial_files_count = len(list(Path(output_folder).glob('*')))
    lst_rag_metadata = []
    if file_type in ['.docx', '.doc']:
        lst_rag_metadata = process_word(chars_per_page, file_name, file_type, output_folder, source_file)
    if file_type == '.pdf':
        lst_rag_metadata = process_pdf(
            chars_per_page, file_name, file_type,
            output_folder, source_file, gen_page_imgs, )
    if file_type in ['.ppt', '.pptx']:
        process_powerpoint(source_file, file_type, output_folder, chars_per_page)
    if file_type in ['.xls', '.xlsx']:
        process_excel(source_file, output_folder, file_name, file_type)
    if file_type == '.ods':
        converted_file_path = convert_office_file(source_file, '.xlsx')
        process_excel(converted_file_path, output_folder, file_name, '.xlsx')
        os.remove(converted_file_path)
    if file_type == '.odp':
        converted_file_path = convert_office_file(source_file, '.pptx')
        process_powerpoint(converted_file_path, '.pptx', output_folder, chars_per_page)
        os.remove(converted_file_path)
    if file_type == '.rtf' or file_type == '.odt':
        converted_file_path = convert_office_file(source_file, '.docx')
        lst_rag_metadata = process_word(chars_per_page, file_name, '.docx', output_folder, converted_file_path)
        os.remove(converted_file_path)
    if file_type == '.txt' or is_text_file(source_file):
        process_txt(chars_per_page, file_name, file_type, output_folder, source_file)
    final_files_count = len(list(Path(output_folder).glob('*')))
    stats['generated'] += (final_files_count - initial_files_count)
    return lst_rag_metadata


def process_txt(chars_per_page, file_name, file_type, output_folder, source_file):
    with open(source_file, 'r', encoding='utf-8') as f:
        text = f.read()
    write_to_file(output_folder, f'{file_name}', text, chars_per_page)


def process_pdf(chars_per_page,
                file_name,
                file_type,
                output_folder,
                source_file,
                gen_page_imgs: bool = False, ) -> List[RagMetadata]:
    doc = fitz.open(source_file)
    lst_rag_metadata = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        text = page.get_text()
        filename = f'{file_name}.{get_page_number_prefix()}{i + 1}'
        write_to_file(output_folder, filename, text, chars_per_page, False)
        # tmp_folders = output_folder.split('/')
        # parent_folder = '/'.join(tmp_folders[:-1])
        lst_rag_metadata.append(RagMetadata(
            name="{filename}.txt".format(filename=filename),
            ori_file_name=f'{file_name}{file_type}',
            gen_page_imgs=gen_page_imgs,
            page_number=(i + 1),
        ))
        # if metadata_dir:
        #     update_download_metadata(metadata_dir,
        #                              f'{file_name}{file_type}',
        #                              "{filename}.txt".format(filename=filename),
        #                              gen_page_imgs)
    return lst_rag_metadata


def process_word(chars_per_page, file_name, file_type, output_folder, source_file) -> List[RagMetadata]:
    pdf_path = convert_office_file(source_file, '.pdf')
    lst_rag_metadata = process_pdf(chars_per_page, file_name, '.pdf', output_folder, pdf_path)
    os.remove(pdf_path)
    return lst_rag_metadata


def split_txts_no_threads(source_files: list,
                          output_folders: list,
                          chars_per_page: int = 2000,
                          force: bool = False,
                          gen_page_imgs: bool = False,
                          metadata_dir: str = None
                          ):
    """
    @param metadata_dir: str, metadata 資料夾路徑，有傳才會更新 metadata
    """
    stats = {'cached': 0, 'generated': 0}
    progress_bar = EmojiProgressBar(len(source_files))  # 初始化進度條
    for index, (source_file, output_folder) in enumerate(zip(source_files, output_folders)):
        try:
            lst_rag_metadata = process_single_file(
                source_file,
                output_folder,
                chars_per_page,
                force,
                stats,
                gen_page_imgs,
            )
            if metadata_dir and len(lst_rag_metadata) > 0:
                update_download_metadata(metadata_dir, lst_rag_metadata)
            progress_bar.update(index + 1)  # 更新進度條
        except Exception as exc:
            log_handler.write_log(f'{source_file} generated an exception: {exc}')
    print(f"split_txts.py, cached files: {stats['cached']}, generated files: {stats['generated']}")


if __name__ == '__main__':
    p = argparse.ArgumentParser(description='Split specified files into pages or parts')
    a = p.add_argument
    a('--source_files', nargs='+', help='Paths of the files to be split')
    a('--output_folders', nargs='+', help='Output folders for the split files')
    a('--chars_per_page', type=int, default=2000)
    a('--force', action='store_true')
    args = p.parse_args()

    split_txts_no_threads(args.source_files, args.output_folders, args.chars_per_page, args.force)
